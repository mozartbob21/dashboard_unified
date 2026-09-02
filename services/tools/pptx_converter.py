from services.shot_stabilizer import stabilize_page
"""HTML -> PPTX конвертер с поддержкой корпоративных шаблонов."""
import base64
import copy
import io
import re
from pathlib import Path

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

TOOLS_DIR = Path(__file__).resolve().parents[2] / "data" / "tools"
TEMPLATES_DIR = TOOLS_DIR / "templates"
OUTPUT_DIR = TOOLS_DIR / "output"
EMBLEM_FILE = TOOLS_DIR / "emblem.png"

GREEN = RGBColor(0x2E, 0x7D, 0x32)
ZEBRA = RGBColor(0xEE, 0xF2, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x3B)


def _clean(text: str) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def _extract_slide(container) -> dict:
    tags = container if isinstance(container, list) else container.find_all(True, recursive=False)
    slide = {"title": "", "subtitle": "", "bullets": [], "texts": [],
             "images": [], "tables": [], "notes": ""}
    for el in tags:
        name = getattr(el, "name", None)
        if not name:
            continue
        if name in ("h1", "h2", "h3") and not slide["title"]:
            slide["title"] = _clean(el.get_text())
        elif name in ("h2", "h3", "h4") and not slide["subtitle"]:
            slide["subtitle"] = _clean(el.get_text())
        elif name in ("ul", "ol"):
            slide["bullets"].extend(_clean(li.get_text()) for li in el.find_all("li"))
        elif name in ("p", "div"):
            if name == "p":
                t = _clean(el.get_text())
            else:
                t = _clean(" ".join(c for c in el.children if isinstance(c, str)))
            if t:
                slide["texts"].append(t)
        elif name == "blockquote":
            t = _clean(el.get_text())
            if t:
                slide["texts"].append("\u00ab" + t + "\u00bb")
        elif name == "table":
            if _is_layout_table(el):
                for cell in el.find_all(["td", "th"]):
                    if not cell.find("table"):
                        t = _clean(cell.get_text())
                        if t:
                            slide["texts"].append(t)
            else:
                rows = [[_clean(td.get_text()) for td in tr.find_all(["td", "th"])]
                        for tr in el.find_all("tr")]
                if rows:
                    slide["tables"].append(rows)
        elif name in ("aside", "div") and "notes" in " ".join(el.get("class", [])):
            slide["notes"] = _clean(el.get_text())
    if not slide["images"]:
        roots = container if isinstance(container, list) else [container]
        imgs = []
        for r in roots:
            imgs += [img.get("src", "") for img in r.find_all("img") if img.get("src")]
        slide["images"] = imgs
    return slide


def _is_layout_table(tbl) -> bool:
    """Таблица-вёрстка: вложенные таблицы, одна колонка или простыня."""
    if tbl.find("table"):
        return True
    rows = tbl.find_all("tr")
    if not rows:
        return False
    widths = {len(tr.find_all(["td", "th"])) for tr in rows}
    if widths <= {1}:
        return True
    if len(rows) == 1 and len(_clean(rows[0].get_text())) > 300:
        return True
    return False


def _unwrap_layout_tables(soup):
    """Разворачивает таблицы-вёрстки в обычные div-блоки (до 6 уровней)."""
    for _ in range(6):
        layouts = [t for t in soup.find_all("table") if _is_layout_table(t)]
        if not layouts:
            break
        for t in layouts:
            div = soup.new_tag("div")
            for cell in t.find_all(["td", "th"]):
                if cell.find("table"):
                    continue
                inner = soup.new_tag("div")
                for ch in list(cell.children):
                    inner.append(ch.extract())
                div.append(inner)
            t.replace_with(div)


def _flatten_unstyled(soup):
    """Снимает div-обёртки без class/id/style и без прямого текста,
    чтобы заголовки и блоки вышли на верхний уровень и разбивка работала."""
    for _ in range(8):
        changed = False
        for div in soup.find_all("div"):
            if div.get("class") or div.get("id") or div.get("style"):
                continue
            direct = "".join(c for c in div.children if isinstance(c, str)).strip()
            if not direct:
                div.unwrap()
                changed = True
        if not changed:
            break


def _extract_presentation_slide(slide):
    """Разбирает один <div class="slide">: заголовок, карточки, таблицы, списки."""
    s = {"title": "", "subtitle": "", "bullets": [], "texts": [],
         "images": [], "tables": [], "notes": ""}
    h = slide.find(["h1", "h2"])
    if h:
        s["title"] = _clean(h.get_text(" "))
    sub = slide.find(class_="subtitle")
    if sub:
        s["subtitle"] = _clean(sub.get_text(" "))
    for card in slide.find_all(class_="stat-card"):
        v = card.find(class_="stat-value")
        l = card.find(class_="stat-label")
        if v:
            s["bullets"].append((_clean(v.get_text(" ")) + " — " +
                                 _clean(l.get_text(" ") if l else "")).rstrip(" —"))
    for box in slide.find_all(class_=["highlight-box", "recommendations-box"]):
        hh = box.find(["h3", "h4"])
        if hh:
            s["bullets"].append(_clean(hh.get_text(" ")) + ":")
        lis = box.find_all("li")
        if lis:
            s["bullets"].extend("• " + _clean(li.get_text(" ")) for li in lis)
        else:
            p = box.find("p")
            if p:
                s["bullets"].append(_clean(p.get_text(" ")))
    for fb in slide.find_all(class_="flow-box"):
        parts = []
        for fi in fb.find_all(class_="flow-item"):
            v = fi.find(class_="flow-value")
            l = fi.find(class_="flow-label")
            chunk = " ".join(x for x in [
                _clean(v.get_text(" ")) if v else "",
                _clean(l.get_text(" ")) if l else ""] if x)
            if chunk:
                parts.append(chunk)
        if parts:
            s["bullets"].append(" → ".join(parts))
    for ul in slide.find_all("ul"):
        if ul.find_parent(class_=["highlight-box", "recommendations-box"]):
            continue
        s["bullets"].extend("• " + _clean(li.get_text(" ")) for li in ul.find_all("li"))
    for t in slide.find_all("table"):
        rows = [[_clean(td.get_text(" ")) for td in tr.find_all(["td", "th"])]
                for tr in t.find_all("tr")]
        rows = [r for r in rows if any(r)]
        if rows:
            s["tables"].append(rows)
    for p in slide.find_all("p"):
        if p.find_parent(class_=["highlight-box", "recommendations-box"]):
            continue
        if "subtitle" in (p.get("class") or []):
            continue
        t = _clean(p.get_text(" "))
        if t:
            s["texts"].append(t)
    for img in slide.find_all("img"):
        if img.get("src"):
            s["images"].append(img["src"])
    return s


def _parse_presentation(soup):
    slides = soup.find_all(class_="slide")
    if not slides:
        return []
    out = [_extract_presentation_slide(x) for x in slides]
    return [x for x in out
            if any([x["title"], x["bullets"], x["texts"], x["tables"], x["images"]])]


def _condense_html(html: str) -> str:
    h = re.sub(r"<(style|script)\b[\s\S]*?</\1>", " ", html, flags=re.I)
    return re.sub(r"\s+", " ", h)



def _extract_slide_div(el):
    """Разбирает один блок <div class="slide"> как один слайд."""
    h = el.find(["h1", "h2", "h3"])
    title = _clean(h.get_text()) if h else ""
    bullets = [_clean(li.get_text()) for li in el.find_all("li")]
    texts = []
    for p in el.find_all("p"):
        t = _clean(p.get_text())
        if t and t != title:
            texts.append(t)
    tables = []
    for t in el.find_all("table"):
        rows = [[_clean(td.get_text()) for td in tr.find_all(["td", "th"])]
                for tr in t.find_all("tr")]
        rows = [r for r in rows if any(r)]
        if rows:
            tables.append(rows)
    images = [img.get("src", "") for img in el.find_all("img") if img.get("src")]
    return {"title": title, "subtitle": "", "bullets": bullets, "texts": texts,
            "images": images, "tables": tables, "notes": ""}


def parse_html_to_slides(html: str) -> list:
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    _flatten_unstyled(soup)

    slide_els = soup.find_all(class_="slide")
    if slide_els:
        slides = [_extract_slide_div(x) for x in slide_els]
        slides = [x for x in slides
                  if any([x["title"], x["bullets"], x["texts"], x["tables"], x["images"]])]
        if slides:
            return slides
    pres = _parse_presentation(soup)
    if pres:
        return pres

    sections = soup.find_all("section")
    if sections:
        slides = [_extract_slide(s) for s in sections]
    else:
        body = soup.body or soup
        split_tag = "h2" if body.find("h2") else "h3"
        groups, cur = [], []
        for el in body.find_all(True, recursive=False):
            if getattr(el, "name", "") == split_tag and cur:
                groups.append(cur)
                cur = []
            cur.append(el)
        if cur:
            groups.append(cur)
        slides = [_extract_slide(g) for g in groups]

    slides = [s for s in slides
              if any([s["title"], s["bullets"], s["texts"], s["images"], s["tables"]])]
    if not slides:
        lines = [ln.strip() for ln in soup.get_text().splitlines() if ln.strip()][:8]
        slides = [{"title": "Презентация", "subtitle": "", "bullets": [],
                   "texts": lines, "images": [],
                   "tables": [], "notes": ""}]
    return slides


def list_templates() -> list:
    TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)
    return sorted(p.name for p in TEMPLATES_DIR.glob("*.pptx"))


def _add_textbox(slide, lines):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(8.8), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = lines[0] if lines else ""
    for ln in lines[1:]:
        tf.add_paragraph().text = ln


def _add_image(slide, src):
    try:
        if src.startswith("data:"):
            data = base64.b64decode(src.split(",", 1)[1])
            slide.shapes.add_picture(io.BytesIO(data), Inches(1.2), Inches(2.2), width=Inches(7.5))
        elif src.startswith("http"):
            import httpx
            r = httpx.get(src, timeout=10)
            if r.status_code == 200:
                slide.shapes.add_picture(io.BytesIO(r.content), Inches(1.2), Inches(2.2), width=Inches(7.5))
        else:
            p = Path(src)
            if p.exists():
                slide.shapes.add_picture(str(p), Inches(1.2), Inches(2.2), width=Inches(7.5))
    except Exception as e:
        print(f"[tools] image error: {e}")


STATUS_FILL = {
    "критический": (RGBColor(0xDC, 0x26, 0x26), WHITE),
    "высокий": (RGBColor(0xFF, 0xC1, 0x07), RGBColor(0x33, 0x33, 0x33)),
    "средний": (RGBColor(0x16, 0xA3, 0x4A), WHITE),
    "низкий": (RGBColor(0x9C, 0xA3, 0xAF), WHITE),
}


def _add_table(slide, rows, left=Inches(0.6), top=Inches(2.3),
               width=Inches(10.5), row_h=Inches(0.5),
               head_size=14, body_size=14, body_bold=False):
    try:
        n_rows, n_cols = len(rows), max(len(r) for r in rows)
        g = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_h * n_rows)
        tbl = g.table
        tbl.first_row = False
        tbl.horz_banding = False
        # ширины колонок по содержимому
        weights = [10] * n_cols
        for r in rows:
            for c in range(min(n_cols, len(r))):
                weights[c] = max(weights[c], min(46, len(str(r[c] or "")) + 4))
        tot = sum(weights)
        for c in range(n_cols):
            tbl.columns[c].width = int(width * weights[c] / tot)
        dense = n_cols >= 6 or n_rows >= 12
        for r_i in range(n_rows):
            tbl.rows[r_i].height = row_h
        for r_i, row in enumerate(rows):
            for c_i in range(n_cols):
                cell = tbl.cell(r_i, c_i)
                cell.text = row[c_i] if c_i < len(row) else ""
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_left = Inches(0.08 if dense else 0.15)
                cell.margin_right = Inches(0.08 if dense else 0.15)
                cell.margin_top = Inches(0.02)
                cell.margin_bottom = Inches(0.02)
                txt = str(cell.text).strip()
                low = txt.lower()
                numeric = bool(re.match(r"^[\d\s.,%()−+-]+$", txt)) and any(ch.isdigit() for ch in txt)
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.RIGHT if (numeric and r_i > 0) else PP_ALIGN.LEFT
                    for run in p.runs:
                        run.font.name = "Calibri"
                        if r_i == 0:
                            run.font.size = Pt(head_size)
                            run.font.bold = True
                            run.font.color.rgb = WHITE
                        else:
                            run.font.size = Pt(body_size)
                            run.font.bold = body_bold or (c_i == 0 and n_cols >= 5)
                            run.font.color.rgb = DARK
                st = STATUS_FILL.get(low)
                if r_i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = GREEN
                elif st:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = st[0]
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.color.rgb = st[1]
                            run.font.bold = True
                elif low.startswith("… ещё"):
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
                            run.font.italic = True
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = ZEBRA if r_i % 2 else WHITE
    except Exception as e:
        print(f"[tools] table error: {e}")


def _load_emblem():
    """Герб из файла data/tools/emblem.png / emblem.jpg (если загружен)."""
    for p in (TOOLS_DIR / "emblem.png", TOOLS_DIR / "emblem.jpg", EMBLEM_FILE):
        try:
            if p.exists():
                return p.read_bytes()
        except Exception:
            continue
    return None


def _find_emblem(prs, index):
    """Герб = самая маленькая картинка титульного слайда."""
    try:
        slide = prs.slides[index]
        pics = [shp for shp in slide.shapes
                if shp.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE)]
        if not pics:
            return None
        small = min(pics, key=lambda sh: sh.width * sh.height)
        blip = small._element.find(".//" + qn("a:blip"))
        rId = blip.get(qn("r:embed"))
        return slide.part.related_part(rId).blob
    except Exception as e:
        print(f"[tools] emblem error: {e}")
        return None


def duplicate_slide(prs, index):
    source = prs.slides[index]
    new = prs.slides.add_slide(prs.slide_layouts[len(prs.slide_layouts) - 1])
    spTree = new.shapes._spTree
    for shp in list(new.shapes):
        spTree.remove(shp._element)
    for shp in source.shapes:
        el = copy.deepcopy(shp._element)
        for b in el.iter(qn("a:blip")):
            for attr in ("r:embed", "r:link"):
                rId = b.get(qn(attr))
                if rId:
                    try:
                        img_part = source.part.related_part(rId)
                        new_rId = new.part.relate_to(img_part, RT.IMAGE)
                        b.set(qn(attr), new_rId)
                    except Exception:
                        pass
        spTree.append(el)
    return new


def _style(tf, size, bold, color, align, v_anchor=None):
    if v_anchor is not None:
        tf.vertical_anchor = v_anchor
    for p in tf.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color


def _fill_slide(slide, s, is_first, prs=None, emblem=None):
    spTree = slide.shapes._spTree
    for shp in list(slide.shapes):
        if shp.has_table:
            spTree.remove(shp._element)
            continue
        if shp.has_text_frame and shp.text_frame.text.strip():
            spTree.remove(shp._element)

    for shp in list(slide.shapes):
        try:
            l, w, h = shp.left, shp.width, shp.height
        except Exception:
            continue
        if shp.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and w < Inches(2):
            spTree.remove(shp._element)
            continue

    if is_first:
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(8.4), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.text = s["title"] or "Презентация"
        _style(tf, 32, True, WHITE, PP_ALIGN.CENTER)
        sub = s["subtitle"] or " ".join(s["texts"])
        if sub:
            tb2 = slide.shapes.add_textbox(Inches(1.2), Inches(4.3), Inches(7.6), Inches(1.0))
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            tf2.text = sub
            _style(tf2, 18, False, WHITE, PP_ALIGN.CENTER)
        for src in s["images"][:2]:
            _add_image(slide, src)
        for rows in s["tables"][:1]:
            _add_table(slide, rows)
        if s["notes"]:
            try:
                slide.notes_slide.notes_text_frame.text = s["notes"]
            except Exception:
                pass
        return

    # ── Контентный слайд ──
    # Правая вертикальная полоса
    right_bar = None
    for shp in slide.shapes:
        try:
            if shp.height > Inches(5) and shp.width < Inches(1.2):
                right_bar = shp
                break
        except Exception:
            continue

    bar_w = (right_bar.left - Inches(0.45)) if right_bar else (
        (prs.slide_width - Inches(1.2)) if prs else Inches(12.1))

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(0.25), bar_w, Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()
    try:
        bar.shadow.inherit = False
    except Exception:
        pass

    # Заголовок белым внутри плашки
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), bar_w - Inches(0.9), Inches(0.85))
    tf = tb.text_frame
    tf.word_wrap = True
    title = s["title"] or ""
    if len(title) > 90:
        title = title[:87] + "…"
    tf.text = title
    _style(tf, 28, True, WHITE, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)

    # Удаляем битые клонированные картинки/группы справа (квадрат с ×)
    for shp in list(slide.shapes):
        try:
            l, w, h = shp.left, shp.width, shp.height
        except Exception:
            continue
        if (shp.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE, MSO_SHAPE_TYPE.GROUP)
                and w < Inches(1.6) and h < Inches(2.2) and l > Inches(7.5)):
            spTree.remove(shp._element)

    # Герб НА правой вертикальной полосе: самый верх, центр полосы
    if emblem and prs is not None:
        try:
            pic = slide.shapes.add_picture(
                io.BytesIO(emblem), Inches(0), Inches(0), height=Inches(1.1))
            bar_center = prs.slide_width - Inches(0.38)
            pic.left = bar_center - pic.width // 2
            pic.top = Inches(0.06)
            print(f"[tools] emblem on bar: left={pic.left / 914400:.2f}in", flush=True)
        except Exception as e:
            print(f"[tools] emblem add error: {e}")


    body = s["bullets"] or s["texts"]
    rows = s["tables"][0] if s["tables"] else None
    top, total_h = 1.35, 5.35
    if body and rows:                      # текст и таблица вместе — делим высоту
        body_h = min(2.4, 0.34 * len(body) + 0.2)
        tbl_top = top + body_h + 0.1
        tbl_h = total_h - body_h - 0.1
    elif body:
        body_h, tbl_top, tbl_h = total_h, None, 0.0
    else:
        body_h, tbl_top, tbl_h = 0.0, top, total_h
    if body:                               # буллеты — сверху
        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(top), bar_w - Inches(0.9), Inches(body_h))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        tf2.text = body[0]
        for ln in body[1:]:
            tf2.add_paragraph().text = ln
        _style(tf2, 18 if len(body) <= 6 else (15 if len(body) <= 10 else 12), False,
               RGBColor(0x33, 0x33, 0x33), PP_ALIGN.LEFT)
        for p in tf2.paragraphs:
            p.space_after = Pt(6)
    for src in s["images"][:2]:
        _add_image(slide, src)
    if rows:                               # таблица — ниже текста
        n_cols = max(len(r) for r in rows)
        if n_cols > 4 or len(rows) > 10:   # широкая — аккуратными строками
            lines = [" | ".join(c for c in r if c) for r in rows]
            cap = 15
            extra = len(lines) - cap
            if extra > 0:
                lines = lines[:cap]
                lines.append("… ещё строк: %d" % extra)
            fs = 14 if len(lines) <= 10 else (11 if len(lines) <= 16 else 9)
            tbb = slide.shapes.add_textbox(Inches(0.6), Inches(tbl_top),
                                           bar_w - Inches(0.9), Inches(tbl_h))
            tfb = tbb.text_frame
            tfb.word_wrap = True
            tfb.text = lines[0] if lines else ""
            for ln in lines[1:]:
                tfb.add_paragraph().text = ln
            _style(tfb, fs, False, RGBColor(0x33, 0x33, 0x33), PP_ALIGN.LEFT)
            for p in tfb.paragraphs:
                p.space_after = Pt(2)
        else:                              # маленькая — настоящей таблицей
            row_h = int(min(Inches(0.9) if len(rows) <= 5 else Inches(0.65),
                            Inches(tbl_h) / max(len(rows), 1)))
            _add_table(slide, rows, left=Inches(0.3), top=Inches(tbl_top),
                       width=bar_w, row_h=row_h,
                       head_size=16, body_size=14, body_bold=True)

    if s["notes"]:
        try:
            slide.notes_slide.notes_text_frame.text = s["notes"]
        except Exception:
            pass


def _build_standard(prs, slides):
    for i, s in enumerate(slides):
        layout_idx = 0 if i == 0 else (1 if s["bullets"] else 5)
        try:
            layout = prs.slide_layouts[layout_idx]
        except IndexError:
            layout = prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title and s["title"]:
            slide.shapes.title.text = s["title"]
        if i == 0 and (s["subtitle"] or s["texts"]):
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    ph.text = s["subtitle"] or " ".join(s["texts"])
                    break
        if i != 0 and (s["bullets"] or s["texts"]):
            body_ph = None
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    body_ph = ph
            lines = s["bullets"] or s["texts"]
            if body_ph is not None:
                tf = body_ph.text_frame
                tf.text = lines[0]
                for ln in lines[1:]:
                    tf.add_paragraph().text = ln
            else:
                _add_textbox(slide, lines)
        for src in s["images"][:2]:
            _add_image(slide, src)
        for rows in s["tables"][:1]:
            _add_table(slide, rows)
        if s["notes"]:
            try:
                slide.notes_slide.notes_text_frame.text = s["notes"]
            except Exception:
                pass


def build_pptx(slides: list, template_name, output_name: str) -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    tpl = TEMPLATES_DIR / template_name if template_name else None
    use_tpl = bool(tpl and tpl.exists())
    print(f"[tools] build_pptx: template={template_name!r}, found={use_tpl}", flush=True)

    if not use_tpl:
        prs = Presentation()
        _build_standard(prs, slides)
    else:
        prs = Presentation(str(tpl))
        n_tpl = len(prs.slides._sldIdLst)
        if n_tpl == 0:
            _build_standard(prs, slides)
        else:
            title_donor = 0
            content_donor = 1 if n_tpl > 2 else 0
            emblem = _load_emblem() or _find_emblem(prs, title_donor)
            for i, s in enumerate(slides):
                donor = title_donor if i == 0 else content_donor
                new_slide = duplicate_slide(prs, donor)
                _fill_slide(new_slide, s, i == 0, prs, emblem)
            sldIdLst = prs.slides._sldIdLst
            orig = list(sldIdLst)[:n_tpl]
            for el in orig: # Убираем срез [:-1], идем по всем
                rId = el.get(qn("r:id"))
                try:
                    prs.part.drop_rel(rId)
                except Exception:
                    pass
                sldIdLst.remove(el)
    out = OUTPUT_DIR / output_name
    prs.save(str(out))
    return out


# ═══════════════ PRO-РАЗБОР СЛОЖНЫХ HTML ═══════════════

def _promote_styled_headings(soup):
    """div/p/span с крупным шрифтом в style -> заголовки h2/h3."""
    for el in soup.find_all(["div", "p", "span"]):
        style = el.get("style") or ""
        text = el.get_text(strip=True)
        if not text or len(text) > 120:
            continue
        fs = re.search(r"font-size:\s*([\d.]+)\s*(px|pt|em)", style)
        fw = re.search(r"font-weight:\s*(bold|[6-9]00)", style)
        size = 0.0
        if fs:
            v = float(fs.group(1))
            size = {"px": v, "pt": v * 1.33, "em": v * 16}[fs.group(2)]
        if size >= 24 or (fw and size >= 16):
            el.name = "h2" if size >= 28 else "h3"
            continue
        # strong/b, занимающие весь блок, = заголовок (табличная верстка, PDF)
        inner = [c for c in el.children if getattr(c, "name", None) in ("strong", "b")]
        if inner and len(text) < 70:
            own = "".join(c for c in el.children if isinstance(c, str)).strip(" :—-·")
            if not own:
                el.name = "h3"


def parse_html_to_slides_pro(html: str) -> list:
    """Умный разбор: стили, pdf2htmlEX (.pf), fallback markdownify."""
    soup = BeautifulSoup(html, "lxml")
    for tag in soup(["script", "style"]):
        tag.decompose()
    _promote_styled_headings(soup)
    _unwrap_layout_tables(soup)

    # pdf2htmlEX: каждая .pf = слайд
    pages = soup.find_all("div", class_="pf")
    if pages:
        slides = []
        for pg in pages:
            seen, bullets = set(), []
            for t in pg.find_all(["h1", "h2", "h3", "p", "li", "div"]):
                # пропускаем обёртки — берём только листовые текстовые блоки
                if t.find(["h1", "h2", "h3", "p", "li", "div"]):
                    continue
                ln = _clean(t.get_text())
                if ln and ln not in seen:
                    seen.add(ln)
                    bullets.append(ln)
            if bullets:
                if len(bullets[0]) <= 90:
                    title, rest = bullets[0], bullets[1:8]
                else:
                    title, rest = "Страница отчёта", bullets[:8]
                slides.append({"title": title, "subtitle": "",
                               "bullets": rest, "texts": [],
                               "images": [], "tables": [], "notes": ""})
        if slides:
            return slides

    slides = parse_html_to_slides(str(soup))
    if len(slides) >= 2:
        return slides

    # fallback: markdownify (любая вложенность)
    try:
        import markdownify
        md = markdownify.markdown(str(soup), heading_style="ATX")
    except Exception:
        return slides

    out, cur = [], None
    for line in md.splitlines():
        s = line.strip()
        if not s:
            continue
        if s.startswith("#") and " " in s[:8]:
            level = len(s) - len(s.lstrip("#"))
            title = s.lstrip("#").strip()
            cur = {"title": title, "subtitle": "", "bullets": [], "texts": [],
                   "images": [], "tables": [], "notes": ""}
            out.append(cur)
            continue
        if cur is None:
            cur = {"title": "Презентация", "subtitle": "", "bullets": [],
                   "texts": [], "images": [], "tables": [], "notes": ""}
            out.append(cur)
        if s.startswith(("- ", "* ")):
            cur["bullets"].append(s[2:])
        elif s.startswith("|"):
            cells = [c.strip() for c in s.strip("|").split("|")]
            if any(cells) and not all(set(c) <= set("-: ") for c in cells):
                cur["tables"].append(cells)
        else:
            cur["texts"].append(s)
    for sl in out:
        if sl["tables"]:
            sl["tables"] = [sl["tables"]]
    return out or slides


def render_slides_images(html: str, stem: str) -> list:
    """Рендер сложного HTML в Chromium: страницы .pf или вертикальные полосы."""
    from playwright.sync_api import sync_playwright
    img_dir = OUTPUT_DIR / f"{stem}_shots"
    img_dir.mkdir(parents=True, exist_ok=True)
    for old in img_dir.glob("*.png"):
        old.unlink()
    paths = []
    with sync_playwright() as pw:
        browser = pw.chromium.launch()
        page = browser.new_page(viewport={"width": 1280, "height": 720})
        page.set_content(html, wait_until="load")
        page.wait_for_timeout(700)
        stabilize_page(page)
        els = page.query_selector_all("div.pf")
        if els:
            for i, el in enumerate(els):
                p = img_dir / f"slide_{i:02d}.png"
                el.screenshot(path=str(p))
                paths.append(p)
        else:
            h = page.evaluate("document.body.scrollHeight")
            y, i = 0, 0
            while y < h and i < 60:
                page.evaluate(f"window.scrollTo(0, {y})")
                page.wait_for_timeout(200)
                p = img_dir / f"slide_{i:02d}.png"
                page.screenshot(path=str(p))
                paths.append(p)
                y += 720
                i += 1
        browser.close()
    return paths


def build_pptx_from_images(images: list, output_name: str) -> Path:
    """Слайды-скриншоты 16:9 во всю площадь."""
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for p in images:
        sl = prs.slides.add_slide(blank)
        pic = sl.shapes.add_picture(str(p), 0, 0)
        sw, sh = prs.slide_width, prs.slide_height
        scale = min(sw / pic.width, sh / pic.height) * 0.96
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        pic.left = int((sw - pic.width) / 2)
        pic.top = int((sh - pic.height) / 2)
    out = OUTPUT_DIR / output_name
    prs.save(str(out))
    return out

# ... (предыдущий код без изменений) ...

# ═══════════════ AI-РАЗБОР (QWEN) ПРЕЗЕНТАЦИЙ ═══════════════

def _slide_source(html: str, limit: int = 16000) -> str:
    """Компактная выжимка исходника: заголовки, списки, таблицы — без разметки."""
    try:
        import markdownify
        md = markdownify.markdown(html, heading_style="ATX")
    except Exception:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html, "html.parser")
        for t in soup(["script", "style"]):
            t.decompose()
        md = soup.get_text("\n")
    md = re.sub(r"[ \t]+\n", "\n", md)
    md = re.sub(r"\n{3,}", "\n\n", md)
    return md.strip()[:limit]

AI_PARSE_PROMPT = (
    "Ты — дизайнер деловых презентаций. Дан текстовый разбор HTML-презентации "
    "(слайды помечены строками «--- СЛАЙД N ---», если автор разметил их сам). "
    "Перерисуй её в строгую структуру слайдов для PowerPoint.\n"
    "Правила:\n"
    "- сохрани все слайды и их порядок, не объединяй и не выкидывай слайды;\n"
    "- title — заголовок слайда, не длиннее 90 символов;\n"
    "- bullets — ключевые факты слайда, не больше 6, каждый не длиннее 110 символов;\n"
    "- таблицы клади в tables: шапка и не более 8 строк; если строк больше — "
    "последней строкой добавь «… и ещё N строк»;\n"
    "- цифры, проценты и названия сохраняй ТОЧНО, ничего не выдумывай;\n"
    "- если слайды не размечены — разбей содержание на 5–10 логических слайдов.\n"
    "Верни СТРОГО JSON-массив без markdown и комментариев:\n"
    '[{"title":"заголовок","bullets":["факт"],"tables":[["Кол1","Кол2"],["1","2"]]}]'
)


def _html_to_brief(html, limit=24000):
    # компактный текстовый разбор: слайды -> заголовки / факты / строки таблиц
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(html, "html.parser")
    for t in soup(["script", "style"]):
        t.decompose()
    blocks = soup.find_all(class_="slide")
    if not blocks:
        blocks = soup.find_all("section")
    parts = []
    if blocks:
        for i, s in enumerate(blocks, 1):
            out, seen = [], set()
            for el in s.find_all(["h1", "h2", "h3", "h4", "p", "li", "tr"]):
                if el.name == "tr":
                    t = " | ".join(" ".join(c.get_text().split()) for c in el.find_all(["td", "th"]))
                else:
                    t = " ".join(el.get_text().split())
                if len(t) < 2 or t in seen:
                    continue
                seen.add(t)
                out.append(t)
                if len(out) >= 40:
                    break
            if out:
                parts.append("--- СЛАЙД %d ---\n%s" % (i, "\n".join(out)))
            if len(parts) >= 30:
                break
        brief = "\n".join(parts)
    else:
        brief = "\n".join(l for l in (" ".join(x.split()) for x in soup.get_text("\n").split("\n")) if len(l) > 1)
    return brief[:limit]


def _parse_ai_json(raw):
    import json as _json
    if not raw:
        return []
    raw = raw.strip()
    raw = re.sub(r"^```[a-zA-Z]*", "", raw).strip()
    if raw.endswith("```"):
        raw = raw[:-3].strip()
    a, b = raw.find("["), raw.rfind("]")
    if a == -1 or b <= a:
        return []
    try:
        data = _json.loads(raw[a:b + 1])
    except Exception:
        return []
    if not isinstance(data, list):
        return []
    out = []
    for it in data:
        if not isinstance(it, dict):
            continue
        title = str(it.get("title") or "").strip()[:90]
        bullets = [str(x).strip()[:140] for x in (it.get("bullets") or []) if str(x).strip()][:8]
        texts = [str(x).strip()[:300] for x in (it.get("texts") or []) if str(x).strip()][:4]
        tables = []
        for t in (it.get("tables") or [])[:2]:
            if isinstance(t, list) and t and all(isinstance(r, list) for r in t):
                rows = [[str(c) for c in r][:8] for r in t][:10]
                if rows:
                    tables.append(rows)
        if title or bullets or texts or tables:
            out.append({"title": title, "subtitle": "", "bullets": bullets,
                        "texts": texts, "images": [], "tables": tables, "notes": ""})
    return out


def _call_ai(msgs, engine):
    from services.summarizer.engine import _qwen_chat
    attempts = []
    if engine:
        attempts.append({"max_tokens": 8000, "engine": engine})
        attempts.append({"max_tokens": 8000, "backend": engine})
    attempts.append({"max_tokens": 8000})
    for kw in attempts:
        try:
            return _qwen_chat(msgs, **kw)
        except TypeError:
            continue
        except Exception as e:
            print("[tools] AI-разбор ошибка: %s" % e, flush=True)
            return None
    print("[tools] AI-разбор: _qwen_chat не принял аргументы", flush=True)
    return None


def _norm_title(t):
    import re as _re
    return _re.sub(r"[^а-яa-z0-9]+", "", str(t or "").lower())


def parse_html_to_slides_ai(html, engine=None):
    brief = _html_to_brief(html)
    msgs = [
        {"role": "system", "content": AI_PARSE_PROMPT},
        {"role": "user", "content": brief},
    ]
    slides = _parse_ai_json(_call_ai(msgs, engine))
    if not slides:
        return slides
    # Модель часто «теряет» таблицы — возвращаем их из структурного разбора
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html, "html.parser")
        st_slides = [_extract_slide_div(el) for el in soup.find_all(class_="slide")]
        if not st_slides:
            st_slides = [_extract_slide_div(soup)]
        for st in st_slides:
            if not st.get("tables"):
                continue
            nt = _norm_title(st["title"])
            tgt = None
            for ai_sl in slides:
                at = _norm_title(ai_sl.get("title"))
                if nt and at and (nt in at or at in nt):
                    tgt = ai_sl
                    break
            if tgt is not None:
                tgt.setdefault("tables", []).extend(st["tables"])
            else:
                slides.append(st)   # табличный слайд, который ИИ совсем выкинул
    except Exception as e:
        print("[tools] merge tables:", e)
    return slides
